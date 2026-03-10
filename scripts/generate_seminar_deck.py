from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = ROOT.parent / "pptx-template.pptx"
OUTPUT_PATH = ROOT / "slides" / "cursor-claude-code-seminar.pptx"

CONTENT_TITLE_BOX = (Emu(1460500), Emu(635000), Emu(10541200), Emu(825600))
PAGE_NUMBER_BOX = (Emu(11366500), Emu(6413500), Emu(635200), Emu(254000))
TITLE_DIVIDER_BOX = (Emu(1460500), Emu(1475000), Emu(10541200), Emu(70000))
BODY_LEFT = Emu(620000)
BODY_TOP = Emu(1720000)
BODY_WIDTH = Emu(10950000)
BODY_HEIGHT = Emu(4320000)
BODY_GAP = Emu(260000)
SECTION_HERO_BOX = (Emu(1080000), Emu(2200000), Emu(10000000), Emu(1700000))
SECTION_NOTE_BOX = (Emu(2000000), Emu(4100000), Emu(8200000), Emu(700000))

TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x42, 0x85, 0xF4)
FONT_BOLD = "Hiragino Kaku Gothic Pro W6"
FONT_REGULAR = "Hiragino Kaku Gothic Pro W3"


def set_paragraph_bullet(paragraph, enabled=False):
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buAutoNum", "a:buBlip", "a:buChar"):
        element = p_pr.find(qn(tag))
        if element is not None:
            p_pr.remove(element)

    if enabled:
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        p_pr.append(bullet)
        paragraph.level = 0
    else:
        bullet_none = OxmlElement("a:buNone")
        p_pr.append(bullet_none)


def add_textbox(slide, box, paragraphs, alignment=None, margins=None):
    left, top, width, height = box
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    if margins is None:
        margins = (0, 0, 0, 0)
    margin_left, margin_right, margin_top, margin_bottom = margins
    frame.margin_left = margin_left
    frame.margin_right = margin_right
    frame.margin_top = margin_top
    frame.margin_bottom = margin_bottom

    for index, paragraph_data in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = alignment
        paragraph.space_after = paragraph_data.get("space_after", Pt(4))
        paragraph.space_before = Pt(0)
        set_paragraph_bullet(paragraph, paragraph_data.get("bullet", False))

        run = paragraph.add_run()
        run.text = paragraph_data["text"]
        run.font.size = paragraph_data["size"]
        run.font.bold = paragraph_data.get("bold", False)
        run.font.name = paragraph_data.get("font_name", FONT_REGULAR)
        run.font.color.rgb = paragraph_data.get("color", TEXT_COLOR)

    return shape


def set_existing_textbox(shape, paragraphs, alignment=None):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, paragraph_data in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = alignment
        paragraph.space_after = paragraph_data.get("space_after", Pt(4))
        paragraph.space_before = Pt(0)
        set_paragraph_bullet(paragraph, paragraph_data.get("bullet", False))

        run = paragraph.add_run()
        run.text = paragraph_data["text"]
        run.font.size = paragraph_data["size"]
        run.font.bold = paragraph_data.get("bold", False)
        run.font.name = paragraph_data.get("font_name", FONT_REGULAR)
        run.font.color.rgb = paragraph_data.get("color", TEXT_COLOR)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides[0]
    set_existing_textbox(
        slide.shapes[0],
        [
            {"text": title, "size": Pt(26), "bold": True, "font_name": FONT_BOLD, "space_after": Pt(12)},
            {"text": subtitle, "size": Pt(15), "font_name": FONT_REGULAR},
        ],
        alignment=PP_ALIGN.CENTER,
    )


def clear_existing_textbox(shape):
    set_existing_textbox(shape, [{"text": "", "size": Pt(1), "color": RGBColor(0xFF, 0xFF, 0xFF)}])


def remove_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]


def add_divider(slide):
    left, top, width, height = TITLE_DIVIDER_BOX
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()


def add_connector(slide, begin, end, color=ACCENT_COLOR, width=Pt(1.8)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin[0], begin[1], end[0], end[1])
    line.line.color.rgb = color
    line.line.width = width
    return line


def add_arrow_connector(slide, begin, end, color=ACCENT_COLOR, width=Pt(1.8)):
    line = add_connector(slide, begin, end, color=color, width=width)
    line.line.end_arrowhead = True
    return line


def add_sequence_participant(slide, x, title, subtitle, *, color):
    header_box = (x, Emu(1820000), Emu(1700000), Emu(680000))
    add_card_background(slide, header_box, RGBColor(0xFC, 0xFD, 0xFF), color)
    add_textbox(
        slide,
        header_box,
        [
            {"text": title, "size": Pt(15), "bold": True, "font_name": FONT_BOLD},
            {"text": subtitle, "size": Pt(11), "font_name": FONT_REGULAR},
        ],
        alignment=PP_ALIGN.CENTER,
        margins=(Emu(70000), Emu(70000), Emu(50000), Emu(30000)),
    )
    center_x = x + Emu(850000)
    lifeline = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, Emu(2500000), center_x, Emu(5600000))
    lifeline.line.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    lifeline.line.width = Pt(1.0)
    return center_x


def add_message(slide, x1, x2, y, text, *, color=ACCENT_COLOR, reverse=False):
    if reverse:
        add_arrow_connector(slide, (x2, y), (x1, y), color=color, width=Pt(1.5))
        left = min(x1, x2) + Emu(30000)
    else:
        add_arrow_connector(slide, (x1, y), (x2, y), color=color, width=Pt(1.5))
        left = min(x1, x2) + Emu(30000)
    width = abs(x2 - x1) - Emu(60000)
    add_textbox(
        slide,
        (left, y - Emu(170000), max(width, Emu(800000)), Emu(220000)),
        [{"text": text, "size": Pt(11), "font_name": FONT_REGULAR}],
        alignment=PP_ALIGN.CENTER,
    )


def add_note_box(slide, box, title, items):
    add_card_background(slide, box, RGBColor(0xFB, 0xFC, 0xFD), RGBColor(0x6B, 0x72, 0x80))
    add_textbox(
        slide,
        (box[0] + Emu(140000), box[1] + Emu(70000), box[2] - Emu(280000), box[3] - Emu(140000)),
        [{"text": title, "size": Pt(14), "bold": True, "font_name": FONT_BOLD, "space_after": Pt(4)}] + normalize_items(items, Pt(11)),
        margins=(Emu(30000), Emu(30000), Emu(30000), Emu(30000)),
    )


def add_sequence_slide(prs, page_number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, CONTENT_TITLE_BOX, [{"text": title, "size": Pt(29), "bold": True, "font_name": FONT_BOLD}])
    add_divider(slide)

    colors = [
        RGBColor(0x42, 0x85, 0xF4),
        RGBColor(0x2F, 0x85, 0x4F),
        RGBColor(0x6B, 0x72, 0x80),
        RGBColor(0xD9, 0x77, 0x06),
        RGBColor(0x0E, 0x74, 0xA8),
        RGBColor(0x11, 0x18, 0x27),
    ]
    headers = [
        (Emu(450000), "ユーザー", "依頼を入力"),
        (Emu(2400000), "入口機能", "Chat / Agent / Commands"),
        (Emu(4350000), "常設ルール", "Rules / AGENTS.md"),
        (Emu(6300000), "進め方整理", "requirements / design / tasks"),
        (Emu(8250000), "補助機能", "Skills / MCP / Multi-agent"),
        (Emu(10200000), "ガードと成果物", "Hooks / docs / tests"),
    ]
    centers = []
    for (x, title_text, subtitle), color in zip(headers, colors):
        centers.append(add_sequence_participant(slide, x, title_text, subtitle, color=color))

    y = Emu(2820000)
    step = Emu(360000)
    add_message(slide, centers[0], centers[1], y, "1. 依頼を入力する")
    y += step
    add_message(slide, centers[1], centers[2], y, "2. 判断基準を参照する", color=RGBColor(0x6B, 0x72, 0x80))
    y += step
    add_message(slide, centers[2], centers[1], y, "3. 守るべき前提を返す", color=RGBColor(0x6B, 0x72, 0x80), reverse=True)
    y += step
    add_message(slide, centers[1], centers[3], y, "4. 要件・設計・タスクを確認する", color=RGBColor(0xD9, 0x77, 0x06))
    y += step
    add_message(slide, centers[3], centers[1], y, "5. 進め方を返す", color=RGBColor(0xD9, 0x77, 0x06), reverse=True)
    y += step
    add_message(slide, centers[1], centers[4], y, "6. 必要なら Skills / MCP / Multi-agent を使う", color=RGBColor(0x0E, 0x74, 0xA8))
    y += step
    add_message(slide, centers[4], centers[1], y, "7. 手順・外部情報・分業結果を返す", color=RGBColor(0x0E, 0x74, 0xA8), reverse=True)
    y += step
    add_message(slide, centers[1], centers[5], y, "8. Hooks でチェックし、コード・文書・差分を更新する", color=RGBColor(0x11, 0x18, 0x27))
    y += step
    add_message(slide, centers[5], centers[0], y, "9. 結果を返す", color=RGBColor(0x2F, 0x85, 0x4F), reverse=True)

    add_note_box(
        slide,
        (Emu(700000), Emu(5780000), Emu(5100000), Emu(500000)),
        "見せたいこと",
        [
            "生成AIが動く中心は入口機能と補助機能",
            "Rules は毎回の判断基準として前提に効く",
            "Hooks は最後の実行時ガードとして効く",
        ],
    )
    add_note_box(
        slide,
        (Emu(6200000), Emu(5780000), Emu(5000000), Emu(500000)),
        "話し方のポイント",
        [
            "Commands は入口、Skills は途中、MCP は外部確認",
            "大きい変更だけ Multi-agent と spec-driven を重く使う",
            "結果は docs / tests / specs へ戻して再利用する",
        ],
    )

    add_textbox(slide, PAGE_NUMBER_BOX, [{"text": str(page_number), "size": Pt(12), "color": ACCENT_COLOR}], alignment=PP_ALIGN.RIGHT)


def add_card_background(slide, box, fill_color, line_color):
    left, top, width, height = box
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.2)
    return shape


def normalize_items(items, default_size):
    paragraphs = []
    for item in items:
        paragraph = dict(item) if isinstance(item, dict) else {"text": item}
        paragraph.setdefault("size", default_size)
        paragraph.setdefault("font_name", FONT_REGULAR)
        paragraph.setdefault("space_after", Pt(5))
        text = paragraph["text"]
        if "bullet" not in paragraph:
            paragraph["bullet"] = not text.startswith(("`", "例:", "画面操作:", "1.", "2.", "3.", "4.", "5."))
        paragraphs.append(paragraph)
    return paragraphs


def add_card(slide, box, heading, items, *, fill_color, accent_rgb, heading_size=Pt(17), default_size=Pt(14)):
    left, top, width, height = box
    add_card_background(slide, box, fill_color, accent_rgb)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + Emu(170000),
        top + Emu(170000),
        Emu(50000),
        Emu(360000),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_rgb
    accent.line.fill.background()

    add_textbox(
        slide,
        (left + Emu(300000), top + Emu(120000), width - Emu(420000), Emu(480000)),
        [{"text": heading, "size": heading_size, "bold": True, "font_name": FONT_BOLD}],
        margins=(Emu(30000), Emu(10000), Emu(10000), 0),
    )
    add_textbox(
        slide,
        (left + Emu(230000), top + Emu(640000), width - Emu(420000), height - Emu(780000)),
        normalize_items(items, default_size),
        margins=(Emu(20000), Emu(20000), Emu(10000), Emu(10000)),
    )


def render_two_card_layout(slide, left_heading, left_items, right_heading, right_items):
    width = (BODY_WIDTH - BODY_GAP) // 2
    left_box = (BODY_LEFT, BODY_TOP, width, BODY_HEIGHT)
    right_box = (BODY_LEFT + width + BODY_GAP, BODY_TOP, width, BODY_HEIGHT)
    add_card(slide, left_box, left_heading, left_items, fill_color=RGBColor(0xFC, 0xFD, 0xFF), accent_rgb=ACCENT_COLOR)
    add_card(slide, right_box, right_heading, right_items, fill_color=RGBColor(0xFC, 0xFC, 0xFD), accent_rgb=RGBColor(0x6B, 0x72, 0x80))


def render_stacked_layout(slide, left_heading, left_items, right_heading, right_items):
    top_height = Emu(2000000)
    bottom_height = BODY_HEIGHT - top_height - BODY_GAP
    top_box = (BODY_LEFT, BODY_TOP, BODY_WIDTH, top_height)
    bottom_box = (BODY_LEFT, BODY_TOP + top_height + BODY_GAP, BODY_WIDTH, bottom_height)
    add_card(slide, top_box, left_heading, left_items, fill_color=RGBColor(0xFC, 0xFD, 0xFF), accent_rgb=ACCENT_COLOR, default_size=Pt(14))
    add_card(slide, bottom_box, right_heading, right_items, fill_color=RGBColor(0xFE, 0xFD, 0xFA), accent_rgb=RGBColor(0xD9, 0x77, 0x06), default_size=Pt(14))


def render_asymmetric_layout(slide, left_heading, left_items, right_heading, right_items):
    left_width = Emu(3600000)
    right_width = BODY_WIDTH - left_width - BODY_GAP
    left_box = (BODY_LEFT, BODY_TOP, left_width, BODY_HEIGHT)
    right_box = (BODY_LEFT + left_width + BODY_GAP, BODY_TOP, right_width, BODY_HEIGHT)
    add_card(slide, left_box, left_heading, left_items, fill_color=RGBColor(0xFB, 0xFC, 0xFD), accent_rgb=RGBColor(0x6B, 0x72, 0x80), default_size=Pt(13))
    add_card(slide, right_box, right_heading, right_items, fill_color=RGBColor(0xFB, 0xFD, 0xFB), accent_rgb=RGBColor(0x2F, 0x85, 0x4F), default_size=Pt(14))


def pick_layout(title):
    stacked_keywords = ("具体的な方法", "進め方", "導入ステップ", "どこに適用", "まとめ")
    asymmetric_keywords = ("実践例", "追加する論点", "レイヤー", "全体像", "ライフサイクル", "デモアプリ")
    if any(keyword in title for keyword in stacked_keywords):
        return "stacked"
    if any(keyword in title for keyword in asymmetric_keywords):
        return "asymmetric"
    return "two_card"


def rewrite_agenda_slide(prs):
    slide = prs.slides[1]
    set_existing_textbox(
        slide.shapes[0],
        [{"text": "今回のアジェンダ", "size": Pt(30), "bold": True, "font_name": FONT_BOLD}],
    )
    clear_existing_textbox(slide.shapes[1])
    set_existing_textbox(
        slide.shapes[2],
        [{"text": "2", "size": Pt(12), "color": ACCENT_COLOR}],
        alignment=PP_ALIGN.RIGHT,
    )
    add_divider(slide)
    render_asymmetric_layout(
        slide,
        "前回の既習前提",
        [
            "Cursorの基本画面、AIチャット、Cmd+K、コンテキスト供給",
            "Eclipseとの併用、プロジェクトルール、自動実行、MCP概論",
            "今回は基本操作ではなく、その先の運用設計を扱う",
        ],
        "今回のゴール",
        [
            "ルール設定、定型コマンド、スキル、MCP、フックの設定と使い分けを理解する",
            "金融開発での具体例と、使わない方がよい場面も判断できるようにする",
            "仕様駆動とマルチエージェントを、チーム運用に接続して考える",
            "導入順と注意点を整理する",
        ],
    )


def add_content_slide(prs, page_number, title, left_heading, left_items, right_heading, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, CONTENT_TITLE_BOX, [{"text": title, "size": Pt(29), "bold": True, "font_name": FONT_BOLD}])
    add_divider(slide)
    layout = pick_layout(title)
    if layout == "stacked":
        render_stacked_layout(slide, left_heading, left_items, right_heading, right_items)
    elif layout == "asymmetric":
        render_asymmetric_layout(slide, left_heading, left_items, right_heading, right_items)
    else:
        render_two_card_layout(slide, left_heading, left_items, right_heading, right_items)
    add_textbox(
        slide,
        PAGE_NUMBER_BOX,
        [{"text": str(page_number), "size": Pt(12), "color": ACCENT_COLOR}],
        alignment=PP_ALIGN.RIGHT,
    )


def render_section_slide(slide, title):
    add_textbox(slide, CONTENT_TITLE_BOX, [{"text": title, "size": Pt(29), "bold": True, "font_name": FONT_BOLD}])
    add_divider(slide)
    add_card_background(slide, SECTION_HERO_BOX, RGBColor(0xF7, 0xFA, 0xFF), ACCENT_COLOR)
    add_textbox(
        slide,
        (Emu(1600000), Emu(2550000), Emu(9000000), Emu(700000)),
        [{"text": title, "size": Pt(30), "bold": True, "font_name": FONT_BOLD, "color": ACCENT_COLOR}],
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        SECTION_NOTE_BOX,
        [{"text": "解決したい問題 → 設定方法 → 実践例 → 注意点 の順に見ます", "size": Pt(18), "font_name": FONT_REGULAR}],
        alignment=PP_ALIGN.CENTER,
    )


def rewrite_intro_section_slide(prs, title):
    slide = prs.slides[2]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            clear_existing_textbox(shape)
    render_section_slide(slide, title)


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_section_slide(slide, title)


def main():
    if not DEFAULT_TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {DEFAULT_TEMPLATE}")

    prs = Presentation(DEFAULT_TEMPLATE)

    add_title_slide(
        prs,
        "Cursor中心で進める生成AI開発の実践",
        "社内向けセミナー | 金融ソリューション開発での再現性と生産性\nルール設定 / 定型コマンド / スキル / 外部接続 / フック / マルチエージェント",
    )
    rewrite_agenda_slide(prs)
    rewrite_intro_section_slide(prs, "1. 今回の論点")
    if len(prs.slides) > 3:
        remove_slide(prs, 3)

    slides = [
        ("sequence", 4, "プロンプト入力後に、どの設定・機能がどう効くか"),
        ("content", 5, "Cursor / Claude Code でAI活用に使う機能の全体像",
         "日々触る機能",
         [
             "AI Chat / Agent: 会話しながら調査、実装、修正を進める",
             "Cmd+K / Inline Edit: その場の編集を素早く依頼する",
             "Background / Cloud Agents: 離席中や並列で作業を進める",
             "Subagents: 調査、実装、レビューの役割を分ける",
         ],
         "運用で効く機能",
         [
             "Rules / AGENTS.md / CLAUDE.md: 判断基準を固定する",
             "Commands: よく使う依頼の入口を短くする",
             "Skills: 長い手順と成果物の型を再利用する",
             "MCP / Hooks / Spec-driven: 外部接続、制御、仕様同期を支える",
         ]),
        ("section", "2. ルール設定"),
        ("content", 7, "ルール設定が解決する問題",
         "ルール設定がない前提",
         [
             "毎回、例外を握りつぶさないことを指示している",
             "人によってログ、命名、テストの基準が揺れる",
             "レビューで同じ形式指摘が何度も発生する",
             "モデルを変えると出力方針がぶれやすい",
         ],
         "ルール設定を入れると",
         [
             "毎回の説明コストを減らせる",
             "AIの出力が現場ルールに近づく",
             "レビューを本質的な論点に寄せやすい",
             "新規参加者もAI経由でお作法を学べる",
         ]),
        ("content", 8, "ルール設定の具体的な方法",
         "作る手順",
         [
             "Explorerで `.cursor/rules/` を作り、用途ごとに `.mdc` を分ける",
             "最初は `coding.mdc` `testing.mdc` `docs-sync.mdc` の3枚で十分",
             "1ファイルは 4〜8 行の短い判断基準から始める",
             "個人の癖はユーザールールへ、repo共通だけgit共有する",
         ],
         "最小例と画面操作",
         [
             {"text": "例: `.cursor/rules/docs-sync.mdc`", "size": Pt(17), "bold": True},
             {"text": "`# Documentation Sync Rules`", "size": Pt(16)},
             {"text": "`- requirements.md / design.md / tasks.md を必要に応じて更新`", "size": Pt(16)},
             {"text": "`- seminar docs と repo 構造をずらさない`", "size": Pt(16)},
             {"text": "画面操作: Explorer → `.cursor` → `rules` → New File → 保存後にAgentで通常依頼し、出力が沿うか確認", "size": Pt(16)},
         ]),
        ("content", 9, "ルール設定の実践例",
         "金融Java開発で入れたい例",
         [
             "Java変更は `@Slf4j` を優先し、例外を握りつぶさない",
             "業務ロジック変更時はテストと文書更新を要求する",
             "DynamoDBや日時型など、ドメイン固有制約を固定する",
             "設計・要件ファイルを更新対象として常に意識させる",
         ],
         "有効な場面",
         [
             "レビューで同じ指摘が繰り返されている",
             "既存コードの様式にAIを合わせたい",
             "チーム外メンバーや新規参加者が増える",
             "複数モデル・複数エージェントで揺れを減らしたい",
         ]),
        ("content", 10, "ルール設定で避けたいこと",
         "使わない方がよい場面",
         [
             "単発依頼や、一度しか使わない特殊手順",
             "詳細な業務手順や長い分岐をルールに詰め込むケース",
             "参考ファイルや例がなく、抽象論だけ並ぶケース",
             "人間向け設計原則をそのまま全部流し込むケース",
         ],
         "注意点",
         [
             "増やしすぎると重要度がぼやける",
             "更新されないルールは誤誘導を生む",
             "TabやInline Editには効かない前提を理解する",
             "最初は3〜7個の高頻度ルールに絞る",
         ]),
        ("section", "3. 定型コマンドとスキル"),
        ("content", 12, "定型コマンドが解決する問題",
         "定型コマンドが向くこと",
         [
             "毎回似た依頼文を書く入口作業",
             "レビュー、プルリクエスト説明、影響確認などの定型開始文",
             "人が会話の主導権を持ちつつAIを動かす用途",
             "短いプロンプトで再現性を持たせたいケース",
         ],
         "定型コマンドが向かないこと",
         [
             "長い手順、分岐、外部ファイルを伴うワークフロー",
             "頻度が低い一発ネタ",
             "成果物の型まで固定したい複雑作業",
             "手順の一貫性が重要な監査・文書作成フロー",
         ]),
        ("content", 13, "定型コマンドの具体的な方法",
         "作る手順",
         [
             "Explorerで `.cursor/commands/` に新規Markdownを作る",
             "ファイル名が `/review-update-impact` のようなコマンド名になる",
             "本文には『何を読んで、何を出すか』だけを書く",
             "長くなったら分割するかスキルへ昇格させる",
         ],
         "最小例と画面操作",
         [
             {"text": "例: `.cursor/commands/review-update-impact.md`", "size": Pt(17), "bold": True},
             {"text": "`# Review Update Impact`", "size": Pt(16)},
             {"text": "`- changed requirements / design / rules / hooks / docs を洗い出す`", "size": Pt(16)},
             {"text": "`- 追加・再実行すべきテストを列挙する`", "size": Pt(16)},
             {"text": "画面操作: Chat入力欄で `/` → コマンド名を選択 → 対象issueや差分を添えて実行", "size": Pt(16)},
         ]),
        ("content", 14, "定型コマンドの実践例",
         "金融開発での例",
         [
             "設計書更新前に `/review-update-impact` を使う",
             "障害票を受けた時に `/bug-triage` で論点を整理する",
             "API仕様変更で `/consumer-check` を起点に影響範囲を見る",
             "定期運用で `/weekly-review` を使って差分を棚卸しする",
         ],
         "注意点",
         [
             "入口に過剰な責務を持たせない",
             "コマンド実行後の判断責任は人が持つ",
             "1コマンドで複数の責務を持たせない",
             "定期的に棚卸しして使われないコマンドは消す",
         ]),
        ("content", 15, "スキルが解決する問題",
         "スキルが向くこと",
         [
             "何度も繰り返す多段階ワークフロー",
             "成果物の型を揃えたい作業",
             "人が変わっても手順を崩したくない作業",
             "要件定義、設計、プルリクエストレビュー、リリース前確認など",
         ],
         "ルール設定、定型コマンドとの違い",
         [
             "ルール設定は短い判断基準",
             "定型コマンドは軽い入口",
             "スキルは長い手順そのもの",
             "長くなったコマンドはスキル化を検討する",
         ]),
        ("content", 16, "スキルの具体的な方法",
         "作る手順",
         [
             "1スキル = 1ディレクトリで、中心に `SKILL.md` を置く",
             "このrepoでは `.cursor/skills/generate-design/` 形式で管理している",
             "長い手順は `scripts/` `references/` `assets/` に分けて同梱する",
             "名前よりも description で『いつ使うか』を明確にする",
         ],
         "最小例と使い方",
         [
             {"text": "例: `.cursor/skills/generate-design/SKILL.md`", "size": Pt(17), "bold": True},
             {"text": "`# Generate Design Skill`", "size": Pt(16)},
             {"text": "`## Purpose` と `## Steps` で目的と実行順を固定する", "size": Pt(16)},
             {"text": "`requirements.md` を読んで `design.md` を更新する、のように入出力を明記する", "size": Pt(16)},
             {"text": "画面操作: Claude Codeでは該当タスクを依頼すると自動起動、必要ならスキル名を明示して呼ぶ", "size": Pt(16)},
         ]),
        ("content", 17, "スキル設計で大事なこと",
         "望ましいスキル",
         [
             "いつ使うかが説明文だけで分かる",
             "成果物の型が明確",
             "前提条件と停止条件が明確",
             "古くなった参照資料が分離されている",
         ],
         "避けたいスキル",
         [
             "対象範囲が広すぎて、何に使うかが曖昧",
             "手順が長いのに成果物の型が曖昧",
             "保守されず、実態とズレた手順を押し付ける",
             "人の判断ポイントが隠れている",
         ]),
        ("content", 18, "スキルの実践例",
         "このリポジトリでの例",
         [
             "`generate-requirements`: 要件定義の型を統一する",
             "`generate-design`: 設計観点を揃える",
             "`implement-feature`: 読む順番、成果物、テスト観点を固定する",
             "`update-docs`: 実装後のドキュメント同期を漏らしにくくする",
         ],
         "有効な場面",
         [
             "人が違っても同じ流れで進めたい時",
             "設計やプルリクエストレビューの品質差を減らしたい時",
             "仕様駆動のやり方をチームで揃えたい時",
             "後から見ても手順が再現できるようにしたい時",
         ]),
        ("content", 19, "定型コマンドとスキルの使い分け",
         "定型コマンドにする",
         [
             "短い入口だけ欲しい",
             "会話の開始文を定型化したい",
             "一人が判断しながら小さく進めたい",
             "複雑な支援ファイルが不要",
         ],
         "スキルにする",
         [
             "多段階で長い手順を再利用したい",
             "成果物の型と前提条件も固定したい",
             "補助スクリプトや参照資料を同梱したい",
             "チーム横断で長く使う運用資産にしたい",
         ]),
        ("section", "4. MCPと外部接続"),
        ("content", 21, "MCPが解決する問題",
         "MCPがないと",
         [
             "課題票、プルリクエスト、設計書、ログ、画面状態を毎回コピペする",
             "AIが古い記憶や不完全な説明で判断しやすい",
             "現在の事実より会話に乗った情報に依存する",
             "外部調査のコストが人に集中する",
         ],
         "MCPを入れると",
         [
             "GitHub、文書、ブラウザ、各種SaaSに直接つなげる",
             "最新の外部事実を前提に判断しやすい",
             "説明と転記のコストを減らせる",
             "AI活用が実務システムに近づく",
         ]),
        ("content", 22, "MCPの具体的な方法",
         "最初の作り方",
         [
             "Settings / Integrations から追加するか `.cursor/mcp.json` を作る",
             "最初は read-only 接続を1つだけ有効化する",
             "ローカル型は `command` / `args`、リモート型は `url` / `headers` で書く",
             "トークンは env や secrets 側で持ち、jsonへ直書きしない",
         ],
         "最小設定例と画面操作",
         [
             {"text": "例: `.cursor/mcp.json`", "size": Pt(17), "bold": True},
             {"text": "`{ \"mcpServers\": { \"docs\": { \"command\": \"python3\", \"args\": [\"<server.py>\"] } } }`", "size": Pt(15)},
             {"text": "`{ \"mcpServers\": { \"remote-docs\": { \"url\": \"<remote-endpoint>\", \"headers\": { \"Authorization\": \"Bearer ${TOKEN}\" } } } }`", "size": Pt(15)},
             {"text": "画面操作: Settingsで接続を有効化 → Chatで『仕様書を見て要約して』のように依頼 → 参照ログを確認", "size": Pt(16)},
         ]),
        ("content", 23, "MCPの実践例 1",
         "開発で分かりやすい例",
         [
             "GitHub MCP: 課題票とプルリクエストの流れを直接読む",
             "ブラウザ用MCP: 画面状態を見て部品単位で確認する",
             "文書系MCP: 公式仕様や一次資料に当たる",
             "監視系MCP: 障害調査の初動を早める",
         ],
         "有効な場面",
         [
             "仕様変更の背景確認",
             "画面の再現と検証",
             "エラー原因の初動調査",
             "設計と実装の突合",
         ]),
        ("content", 24, "MCPの実践例 2",
         "金融開発での使いどころ",
         [
             "顧客向け提案書テンプレートや社内標準文書の参照",
             "課題票 / 議事録 / 仕様書を横断した変更整理",
             "画面、ログ、仕様の3点を見ながら障害解析",
             "監査用の説明資料作成時に一次情報へすぐ戻れる",
         ],
         "使わない方がよい場面",
         [
             "ローカルコードだけで足りる時",
             "権限の強い接続を雑に配布する時",
             "監査や権限境界が未整備なまま導入する時",
             "更新権限が不要なのに書き込みを許している時",
         ]),
        ("content", 25, "MCPの注意点",
         "最低限決めること",
         [
             "誰が何のためにそのMCPを使うか",
             "読み取り / 書き込みの境界",
             "認証方式と認証切れ時の扱い",
             "ログ・監査・権限棚卸しの方法",
         ],
         "失敗パターン",
         [
             "便利だからと何でもつなぐ",
             "権限が肥大化して誰も把握できない",
             "利用目的が曖昧で運用ルールがない",
             "接続先が増えたのに文書が更新されない",
         ]),
        ("section", "5. フック"),
        ("content", 27, "フックが解決する問題",
         "フックがないと",
         [
             "危険操作の確認が人任せになる",
             "毎回同じ整形、静的検査、文書確認を手でやる",
             "チームルールを破っても後追いでしか気づけない",
             "AI活用を広げるほど事故率が上がりやすい",
         ],
         "フックを入れると",
         [
             "エージェント処理の前後にガードレールを置ける",
             "危険コマンドやMCP操作を遮断・警告できる",
             "編集後の自動チェックを差し込める",
             "利用ログや監査情報も取りやすい",
         ]),
        ("content", 28, "フックの具体的な方法",
         "作る手順",
         [
             "先に `hooks/` 配下へシェルスクリプトを置き、実行権限を付ける",
             "次に `.cursor/hooks.json` でイベントとスクリプトを結び付ける",
             "最初は警告だけ出す軽いチェックから始める",
             "重い test は保存ごとではなく、明示実行かコミット前へ寄せる",
         ],
         "最小設定例と画面操作",
         [
             {"text": "例: `.cursor/hooks.json`", "size": Pt(17), "bold": True},
             {"text": "`{ \"afterFileEdit\": [{ \"command\": \"hooks/docs-sync-check.sh\" }] }`", "size": Pt(15)},
             {"text": "よく使うイベント: `beforeShellExecution` `afterFileEdit` `beforeMCPExecution`", "size": Pt(16)},
             {"text": "画面操作: 保存やshell実行の後にHookログを確認し、止める設計は慣れてから入れる", "size": Pt(16)},
         ]),
        ("content", 29, "フックの実践例",
         "入れやすい例",
         [
             "危険な shell を拒否する",
             "編集後に整形や静的検査を実行する",
             "仕様文書が古そうなら警告を出す",
             "MCPの書き込み操作の前に理由入力を求める",
         ],
         "金融開発での例",
         [
             "機密文字列の混入チェック",
             "監査対象ディレクトリへの変更ログ収集",
             "設計書未更新時の警告",
             "DB更新系コマンドの制限",
         ]),
        ("content", 30, "フックで避けたいこと",
         "使わない方がよい設計",
         [
             "重くて遅いフックを大量に入れる",
             "失敗理由が分からないブラックボックス制御",
             "曖昧な判定で頻繁に作業を止める",
             "ローカル権限で危険なスクリプトを実行する",
         ],
         "注意点",
         [
             "fail-open / deny 条件を理解して設計する",
             "まず観測だけ行い、後からブロックへ強化する",
             "フック自体のコードレビューが必要",
             "『守りたいもの』を明確にしてから入れる",
         ]),
        ("section", "6. 複数エージェントと仕様駆動"),
        ("content", 32, "マルチエージェントの使いどころ",
         "分業に向く仕事",
         [
             "調査担当: 調査や関連箇所探索",
             "実装担当: 独立タスクの実装",
             "レビュー担当: リスク、欠落、テスト観点確認",
             "画面確認担当: 再現確認や見た目検証",
         ],
         "分けない方がよい仕事",
         [
             "小さい修正",
             "密結合で衝突しやすい変更",
             "役割分担が曖昧なままの並列化",
             "最後の統合判断をAIだけに任せるケース",
         ]),
        ("content", 33, "クラウド実行 / worktree / 並列実行の考え方",
         "Cursor公式で押さえる点",
         [
             "クラウドエージェントはクラウド上で並列に走る",
             "別ブランチや分離環境でビルドとテストができる",
             "Slack / GitHub / Web からも起動できる",
             "ローカルを離れても進められる",
         ],
         "使いどころ",
         [
             "ドキュメント更新やテスト追加のような積み残し",
             "大きめのリファクタの分割実行",
             "レビュー観点を複数モデルで比較したい時",
             "夜間や離席中に回したい作業",
         ]),
        ("content", 34, "仕様駆動の進め方",
         "進め方",
         [
             "`requirements`: 何を解くか",
             "`design`: どう作るか",
             "`tasks`: どう分割するか",
             "実装後は文書とテストを同期する",
         ],
         "いつ重く回すか",
         [
             "新機能",
             "仕様が曖昧な変更",
             "影響範囲が広い変更",
             "金融系で説明責任が重い変更",
         ]),
        ("content", 35, "デモアプリでどこに適用するか",
         "Portfolio Copilot への対応",
         [
             "ルール設定: Java / ログ出力 / テスト / 文書同期",
             "定型コマンド: 要件整理、設計更新、影響確認",
             "スキル: 仕様駆動の一連手順",
             "MCP: ブラウザ / 文書 / GitHub / API仕様確認",
         ],
         "実演の画面操作",
         [
             "1. `./gradlew bootRun` 後、`/portfolios` を開いて一覧を見せる",
             "2. 1件選び、相談文を入力して会話履歴が残ることを見せる",
             "3. 提案書ドラフト生成を押し、議論内容が反映されることを確認する",
             "4. Cursorへ戻り `docs/specs` → `.cursor` → `hooks/` → `src/test` の順に開く",
         ]),
        ("content", 36, "機能追加時に整合性をどう保つか",
         "更新フロー",
         [
             "まず requirements の差分を決める",
             "必要に応じて design と tasks を更新する",
             "実装後にテストと文書の差分を確認する",
             "必要ならルール設定、スキル、フックも更新する",
         ],
         "よくある失敗",
         [
             "コードだけ変わり、仕様書が古いまま残る",
             "手順変更後もスキルが古く、誤った誘導を続ける",
             "フックが重くなり、誰も使わなくなる",
             "MCP権限が肥大化し、監査しづらくなる",
         ]),
        ("content", 37, "Cursor中心だが Claude Code にも通じる",
         "Cursorで見る対応関係",
         [
             "ルール設定 = 常設の判断基準",
             "定型コマンド = 軽い入口",
             "スキル = 長手順の再利用",
             "MCP = 外部接続",
             "フック / マルチエージェント = 制御と分業",
         ],
         "Claude Codeでの近い考え方",
         [
             "スキル / 独自コマンド は同様に使える",
             "フック と サブエージェント で制御と分業ができる",
             "MCPは接続、スキルは使い方の整理に向く",
             "違いより、運用設計を共通化する方が重要",
         ]),
        ("content", 38, "導入ステップ: 明日から始める順番",
         "最初の4週間",
         [
             "1週目: 横断ルールを3つだけ作る",
             "2週目: 繰り返す依頼を2つコマンド化する",
             "3週目: 長い定型作業を1つスキル化する",
             "4週目: MCP / フック / マルチエージェント を限定導入する",
         ],
         "社内展開のポイント",
         [
             "まず1チームで回し、効いた型だけ残す",
             "効果指標は速度だけでなく手戻り率も見る",
             "金融開発では説明可能性を成果指標に入れる",
             "前回資料の基本編と今回の運用編をセットで使う",
         ]),
        ("content", 39, "まとめ",
         "確認したい点",
         [
             "便利機能は知るだけでは成果にならない",
             "ルール設定、スキル、MCP、フックは役割が違う",
             "速く書くより、継続的に運用できる開発フローを作る",
             "Cursor中心で始め、Claude Codeにも広げられる",
         ],
         "次の一歩",
         [
             "まず自チームの頻出レビュー指摘を3つルール化する",
             "次に2つの定型依頼をコマンド化する",
             "長手順を1つスキル化する",
             "その後にMCPとフックで実務接続を強める",
         ]),
    ]

    for item in slides:
        if item[0] == "section":
            add_section_slide(prs, item[1])
        elif item[0] == "sequence":
            _, page_number, title = item
            add_sequence_slide(prs, page_number, title)
        else:
            _, page_number, title, left_heading, left_items, right_heading, right_items = item
            add_content_slide(prs, page_number, title, left_heading, left_items, right_heading, right_items)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"generated: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
