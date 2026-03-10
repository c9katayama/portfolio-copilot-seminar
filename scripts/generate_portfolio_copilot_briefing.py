from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "docs" / "slides" / "portfolio-copilot-briefing.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(0x16, 0x32, 0x4F)
NAVY_DEEP = RGBColor(0x0E, 0x22, 0x38)
SLATE = RGBColor(0x5B, 0x7C, 0x99)
INK = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x52, 0x65, 0x78)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xEA, 0xF0, 0xF6)
PALE_NAVY = RGBColor(0xE3, 0xEB, 0xF3)
SOFT_GREEN = RGBColor(0xE9, 0xF4, 0xEF)
SOFT_AMBER = RGBColor(0xF8, 0xF0, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC8, 0xD3, 0xDF)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = "Cursor"
    prs.core_properties.title = "Portfolio Copilot 説明資料"
    prs.core_properties.subject = "Portfolio Copilot application briefing"
    prs.core_properties.language = "ja-JP"
    return prs


def full_background(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs_w(), prs_h())
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def prs_w():
    return Inches(SLIDE_W)


def prs_h():
    return Inches(SLIDE_H)


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    return shape


def add_plain_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    *,
    margin=0.08,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign

    for idx, item in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item["text"]
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(item.get("space_after", 0))
        paragraph.space_before = Pt(0)
        font = paragraph.runs[0].font
        font.name = item.get("font_name", FONT_BODY)
        font.size = Pt(item.get("size", 16))
        font.bold = item.get("bold", False)
        font.color.rgb = item.get("color", INK)
    return box


def add_badge(slide, x, y, w, h, text, *, fill=PALE_NAVY, color=NAVY):
    add_rect(slide, x, y, w, h, fill)
    add_textbox(
        slide,
        x + 0.08,
        y + 0.03,
        w - 0.16,
        h - 0.06,
        [{"text": text, "size": 11, "bold": True, "color": color, "align": PP_ALIGN.CENTER}],
        margin=0.0,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_number_circle(slide, x, y, diameter, number, *, fill=NAVY, color=WHITE, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_textbox(
        slide,
        x,
        y,
        diameter,
        diameter,
        [{"text": str(number), "size": font_size, "bold": True, "color": color, "align": PP_ALIGN.CENTER}],
        margin=0.0,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_title(slide, title, subtitle=None, *, dark=False, page=None):
    title_color = WHITE if dark else INK
    sub_color = PALE_BLUE if dark else MUTED
    add_textbox(
        slide,
        0.7,
        0.5,
        9.7,
        1.2,
        [{"text": title, "size": 28, "bold": True, "font_name": FONT_HEAD, "color": title_color}],
        margin=0.0,
    )
    if subtitle:
        add_textbox(
            slide,
            0.72,
            1.3,
            9.6,
            0.65,
            [{"text": subtitle, "size": 12, "color": sub_color}],
            margin=0.0,
        )
    if page is not None:
        add_badge(slide, 11.95, 0.45, 0.8, 0.34, str(page), fill=PALE_BLUE if dark else PALE_NAVY, color=NAVY)


def add_footer_note(slide, text, *, dark=False):
    fill = RGBColor(0x20, 0x3D, 0x5B) if dark else PALE_BLUE
    color = WHITE if dark else MUTED
    add_plain_rect(slide, 0.7, 6.78, 11.95, 0.38, fill)
    add_textbox(
        slide,
        0.88,
        6.81,
        11.55,
        0.26,
        [{"text": text, "size": 10, "color": color}],
        margin=0.0,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_metric_card(slide, x, y, w, h, number, label, detail):
    add_rect(slide, x, y, w, h, CARD_BG, line=PALE_NAVY)
    add_textbox(
        slide,
        x + 0.18,
        y + 0.15,
        w - 0.36,
        0.6,
        [{"text": number, "size": 24, "bold": True, "font_name": FONT_HEAD, "color": NAVY}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 0.84,
        w - 0.36,
        0.38,
        [{"text": label, "size": 10.5, "bold": True, "color": MUTED}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 1.22,
        w - 0.36,
        0.56,
        [{"text": detail, "size": 9.5, "color": MUTED}],
        margin=0.0,
    )


def add_feature_card(slide, x, y, w, h, number, title, body, tint):
    add_rect(slide, x, y, w, h, CARD_BG, line=LINE)
    add_number_circle(slide, x + 0.18, y + 0.18, 0.46, number, fill=tint, font_size=14)
    add_textbox(
        slide,
        x + 0.75,
        y + 0.18,
        w - 0.95,
        0.45,
        [{"text": title, "size": 16, "bold": True, "font_name": FONT_HEAD}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 0.82,
        w - 0.36,
        h - 1.0,
        [{"text": body, "size": 11, "color": MUTED}],
        margin=0.0,
    )


def add_step_card(slide, x, y, w, h, number, title, body, fill):
    add_rect(slide, x, y, w, h, fill)
    add_number_circle(slide, x + 0.18, y + 0.18, 0.44, number, fill=NAVY, font_size=13)
    add_textbox(
        slide,
        x + 0.74,
        y + 0.18,
        w - 0.92,
        0.4,
        [{"text": title, "size": 14, "bold": True, "font_name": FONT_HEAD}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 0.72,
        w - 0.36,
        h - 0.88,
        [{"text": body, "size": 11, "color": MUTED}],
        margin=0.0,
    )


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = SLATE
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def add_arch_node(slide, x, y, w, h, number, title, subtitle, fill):
    add_rect(slide, x, y, w, h, fill, line=LINE)
    add_number_circle(slide, x + 0.16, y + 0.16, 0.42, number, fill=NAVY, font_size=13)
    add_textbox(
        slide,
        x + 0.68,
        y + 0.14,
        w - 0.84,
        0.4,
        [{"text": title, "size": 11.5, "bold": True, "font_name": FONT_HEAD}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 0.66,
        w - 0.36,
        h - 0.82,
        [{"text": line, "size": 8.6, "color": MUTED} for line in subtitle.split("\n")],
        margin=0.0,
    )


def add_stack_card(slide, x, y, w, h, title, lines, fill):
    add_rect(slide, x, y, w, h, fill, line=LINE)
    add_textbox(
        slide,
        x + 0.18,
        y + 0.16,
        w - 0.36,
        0.34,
        [{"text": title, "size": 12, "bold": True, "font_name": FONT_HEAD}],
        margin=0.0,
    )
    add_textbox(
        slide,
        x + 0.18,
        y + 0.55,
        w - 0.36,
        h - 0.7,
        [{"text": line, "size": 10.4, "color": MUTED} for line in lines],
        margin=0.0,
    )


def build_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, NAVY_DEEP)

    add_plain_rect(slide, 8.75, 0.0, 4.58, 7.5, NAVY)
    add_plain_rect(slide, 8.1, 0.85, 4.9, 5.7, RGBColor(0x21, 0x45, 0x67))
    add_plain_rect(slide, 8.55, 1.25, 4.0, 1.0, SLATE)
    add_plain_rect(slide, 8.55, 2.55, 3.2, 0.85, PALE_BLUE)
    add_plain_rect(slide, 8.55, 3.72, 4.15, 1.25, RGBColor(0xD9, 0xE6, 0xF2))
    add_plain_rect(slide, 8.55, 5.32, 2.75, 0.62, SOFT_GREEN)

    add_badge(slide, 0.72, 0.75, 1.8, 0.38, "アプリ説明資料", fill=RGBColor(0x23, 0x4B, 0x70), color=WHITE)
    add_title(
        slide,
        "Portfolio Copilot",
        "ポートフォリオ確認、AI対話、提案書ドラフト生成の流れを整理した資料",
        dark=True,
    )
    add_textbox(
        slide,
        0.72,
        2.15,
        6.8,
        1.4,
        [
            {
                "text": "金融業務を題材にしたローカル実行アプリとして、利用者視点の操作と構成を順に確認できる。",
                "size": 16,
                "color": WHITE,
            }
        ],
        margin=0.0,
    )

    add_metric_card(slide, 0.72, 4.2, 2.15, 1.82, "3", "サンプルポートフォリオ", "SQLite へ初期登録")
    add_metric_card(slide, 3.08, 4.2, 2.15, 1.82, "2", "AI プロバイダ", "OpenAI 既定 / Anthropic 切替")
    add_metric_card(slide, 5.44, 4.2, 2.15, 1.82, "4", "説明テーマ", "機能、操作、構成、環境")
    add_footer_note(slide, "README、仕様書、実装コードをもとに事実ベースで構成している。", dark=True)


def build_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, LIGHT_BG)
    add_title(slide, "1. アプリ概要", "利用場面と、このアプリで短時間に把握できること", page=1)

    add_rect(slide, 0.7, 1.6, 4.15, 4.6, CARD_BG, line=PALE_NAVY)
    add_badge(slide, 0.95, 1.9, 1.55, 0.34, "利用シーン")
    add_textbox(
        slide,
        0.95,
        2.38,
        3.55,
        1.35,
        [
            {"text": "金融機関での提案準備や説明の流れを題材にしたローカル実行のサンプルアプリ", "size": 18, "bold": True, "font_name": FONT_HEAD},
            {"text": "実運用の確定判断ではなく、検討用ドラフトの整理と説明を目的とする", "size": 11, "color": MUTED},
        ],
        margin=0.0,
    )
    add_textbox(
        slide,
        0.95,
        4.0,
        3.55,
        1.5,
        [
            {"text": "画面で確認できる内容", "size": 12, "bold": True, "color": NAVY},
            {"text": "1. ポートフォリオ概要と保有比率", "size": 11, "color": MUTED},
            {"text": "2. AI との相談履歴", "size": 11, "color": MUTED},
            {"text": "3. 最新の提案書ドラフト", "size": 11, "color": MUTED},
        ],
        margin=0.0,
    )

    add_feature_card(
        slide,
        5.15,
        1.7,
        3.55,
        1.9,
        1,
        "利用者が見る中心画面",
        "一覧から対象ポートフォリオを選ぶと、目的、リスク、投資金額、保有明細を同じ画面で確認できる。",
        NAVY,
    )
    add_feature_card(
        slide,
        8.98,
        1.7,
        3.65,
        1.9,
        2,
        "相談の流れが切れにくい",
        "AI への相談結果は保存され、画面を再表示しても履歴として残る。後続の提案書生成にも流用できる。",
        SLATE,
    )
    add_feature_card(
        slide,
        5.15,
        4.05,
        3.55,
        1.9,
        3,
        "ローカルで起動しやすい",
        "Spring Boot と SQLite を使い、クラウド配備なしでデモ可能。環境差分を小さくしやすい。",
        RGBColor(0x3E, 0x6A, 0x52),
    )
    add_feature_card(
        slide,
        8.98,
        4.05,
        3.65,
        1.9,
        4,
        "構成を説明しやすい",
        "Web、サービス、SQLite、AI API の役割が分かれており、処理の流れを追いやすい。",
        RGBColor(0x9B, 0x6B, 0x29),
    )
    add_footer_note(slide, "README、仕様書、主要実装を参照して説明内容を構成している。")


def build_features_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, WHITE)
    add_title(slide, "2. 主な機能", "利用者目線で見たときに何ができるか", page=2)

    add_feature_card(
        slide,
        0.7,
        1.7,
        3.0,
        2.0,
        1,
        "ポートフォリオ閲覧",
        "一覧からサンプルデータを選択し、目的、顧客区分、リスク属性、投資金額を把握できる。",
        NAVY,
    )
    add_feature_card(
        slide,
        3.98,
        1.7,
        3.0,
        2.0,
        2,
        "保有比率の比較表示",
        "Allocation Overview と一覧表の両方で、保有銘柄、比率、市場価値を確認できる。",
        SLATE,
    )
    add_feature_card(
        slide,
        7.26,
        1.7,
        3.0,
        2.0,
        3,
        "AI 対話支援",
        "相談内容を送ると、ポートフォリオ文脈と過去履歴を使って論点整理の回答を返す。",
        RGBColor(0x3E, 0x6A, 0x52),
    )
    add_feature_card(
        slide,
        10.54,
        1.7,
        2.1,
        2.0,
        4,
        "提案書ドラフト生成",
        "会話履歴を踏まえ、最新ドラフトを画面内に表示する。",
        RGBColor(0x9B, 0x6B, 0x29),
    )

    add_rect(slide, 0.7, 4.15, 5.7, 1.85, PALE_BLUE)
    add_textbox(
        slide,
        0.95,
        4.42,
        5.2,
        1.35,
        [
            {"text": "ユーザーにとっての見え方", "size": 13, "bold": True, "font_name": FONT_HEAD, "color": NAVY},
            {"text": "単独画面で、確認、相談、ドラフト生成を続けて実行できる。業務説明の途中で画面を行き来しなくてよい。", "size": 12, "color": MUTED},
        ],
        margin=0.0,
    )

    add_rect(slide, 6.7, 4.15, 5.95, 1.85, SOFT_GREEN)
    add_textbox(
        slide,
        6.95,
        4.42,
        5.4,
        1.35,
        [
            {"text": "注意点", "size": 13, "bold": True, "font_name": FONT_HEAD, "color": RGBColor(0x3E, 0x6A, 0x52)},
            {"text": "AI 出力は内部検討用のドラフト支援として扱う。設定不足や API 失敗時は明示的なエラー応答を返す設計になっている。", "size": 12, "color": MUTED},
        ],
        margin=0.0,
    )
    add_footer_note(slide, "要件書の 2. Portfolio Management、3. AI Discussion Support、4. Proposal Draft Generation に対応。")


def build_operation_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, LIGHT_BG)
    add_title(slide, "3. 操作方法", "初見の利用者が追いやすい順序で整理した操作フロー", page=3)

    y = 2.0
    cards = [
        (0.72, 1, "対象を選ぶ", "一覧からポートフォリオを開く。既定では最初の対象が表示される。", PALE_BLUE),
        (3.72, 2, "配分を確認", "目的、リスク、投資金額、保有銘柄と配分を確認する。", CARD_BG),
        (6.72, 3, "AIに相談", "相談内容を送信する。正常時は再読込後に履歴へ反映される。", SOFT_GREEN),
        (9.72, 4, "ドラフトを生成", "提案書ドラフトを生成し、会話履歴とあわせて見直す。", SOFT_AMBER),
    ]
    for x, number, title, body, fill in cards:
        add_step_card(slide, x, y, 2.45, 2.55, number, title, body, fill)

    for start_x in (3.17, 6.17, 9.17):
        add_arrow(slide, start_x, 3.25, start_x + 0.35, 3.25)

    add_rect(slide, 0.85, 5.15, 12.0, 0.9, CARD_BG, line=LINE)
    add_textbox(
        slide,
        1.05,
        5.38,
        11.6,
        0.4,
        [{"text": "補足: API エラーや設定不足がある場合は status 領域にメッセージを出し、黙って既定値に置き換えることはしない。", "size": 11, "color": MUTED}],
        margin=0.0,
    )
    add_footer_note(slide, "画面操作は `GET /portfolios` と `POST /api/portfolios/{id}/chat` / `proposal` を起点にしている。")


def build_architecture_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, WHITE)
    add_title(slide, "4. アーキテクチャ", "UI からデータ保存、AI 連携までの流れ", page=4)

    add_arch_node(slide, 0.8, 1.8, 2.0, 1.1, 1, "Browser UI", "Thymeleaf\nVanilla JS", PALE_BLUE)
    add_arch_node(slide, 3.1, 1.8, 2.0, 1.1, 2, "Web Controller", "PortfolioController\nAiController", PALE_NAVY)
    add_arch_node(slide, 5.4, 1.8, 2.2, 1.1, 3, "Application Service", "Portfolio / Discussion /\nProposal Service", SOFT_GREEN)
    add_arch_node(slide, 7.95, 1.2, 2.15, 1.1, 4, "JDBC Repo", "Portfolio / Holding\nConversation / Proposal", CARD_BG)
    add_arch_node(slide, 7.95, 3.0, 2.15, 1.0, 5, "SQLite", "portfolio-copilot.db\nschema.sql / data.sql", SOFT_AMBER)
    add_arch_node(slide, 10.45, 1.8, 2.0, 1.1, 6, "AI API Client", "OpenAI default\nAnthropic optional", PALE_BLUE)
    add_arch_node(slide, 5.25, 4.1, 2.5, 0.95, 7, "Prompt / Config", "PromptBuilder\nAI provider settings", CARD_BG)

    add_arrow(slide, 2.8, 2.35, 3.1, 2.35)
    add_arrow(slide, 5.1, 2.35, 5.4, 2.35)
    add_arrow(slide, 7.6, 2.35, 8.0, 1.75)
    add_arrow(slide, 7.6, 2.35, 10.45, 2.35)
    add_arrow(slide, 9.05, 2.3, 9.05, 3.0)
    add_arrow(slide, 6.5, 4.1, 6.5, 2.95)

    add_rect(slide, 0.8, 5.28, 3.7, 0.85, PALE_BLUE)
    add_textbox(
        slide,
        1.02,
        5.48,
        3.28,
        0.4,
        [
            {"text": "ポイント 1", "size": 12, "bold": True, "font_name": FONT_HEAD, "color": NAVY},
            {"text": "UI とサービスを分け、永続化と AI 呼び出しをサービス層で集約する。", "size": 10.5, "color": MUTED},
        ],
        margin=0.0,
    )
    add_rect(slide, 4.8, 5.28, 3.7, 0.85, SOFT_GREEN)
    add_textbox(
        slide,
        5.02,
        5.48,
        3.28,
        0.4,
        [
            {"text": "ポイント 2", "size": 12, "bold": True, "font_name": FONT_HEAD, "color": RGBColor(0x3E, 0x6A, 0x52)},
            {"text": "履歴とドラフトは SQLite に保存され、画面再表示でも確認できる。", "size": 10.5, "color": MUTED},
        ],
        margin=0.0,
    )
    add_rect(slide, 8.8, 5.28, 3.65, 0.85, SOFT_AMBER)
    add_textbox(
        slide,
        9.02,
        5.48,
        3.23,
        0.4,
        [
            {"text": "ポイント 3", "size": 12, "bold": True, "font_name": FONT_HEAD, "color": RGBColor(0x9B, 0x6B, 0x29)},
            {"text": "プロンプト生成と AI 接続設定を分け、応答生成の前提を把握しやすくしている。", "size": 10.5, "color": MUTED},
        ],
        margin=0.0,
    )
    add_footer_note(slide, "図中のノードはすべて番号付き。design.md の構成をスライド向けに整理している。")


def build_stack_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, LIGHT_BG)
    add_title(slide, "5. 技術スタック", "役割ごとに分けて把握できる構成", page=5)

    add_stack_card(slide, 0.75, 1.65, 3.9, 1.4, "UI と操作", ["Thymeleaf", "Vanilla JS", "単一画面で確認"], CARD_BG)
    add_stack_card(slide, 0.75, 3.2, 3.9, 1.4, "アプリ本体", ["Spring Boot 3.5", "Java 21", "Controller・Service 分離"], CARD_BG)
    add_stack_card(slide, 0.75, 4.75, 3.9, 1.4, "データ保存", ["SQLite", "JdbcClient", "schema.sql・data.sql"], CARD_BG)

    add_stack_card(slide, 4.95, 1.65, 3.7, 1.4, "AI 連携", ["OpenAI 既定", "Anthropic 切替可", "HttpClient 利用"], PALE_BLUE)
    add_stack_card(slide, 4.95, 3.2, 3.7, 1.4, "品質維持", ["Spring Boot Test", "Spotless", "明示的なエラー処理"], SOFT_GREEN)
    add_stack_card(slide, 4.95, 4.75, 3.7, 1.4, "プロンプトと設定", ["PromptBuilder", "OpenAI / Anthropic 切替", "API キー必須"], PALE_NAVY)

    add_rect(slide, 8.95, 1.7, 3.7, 4.15, NAVY)
    add_textbox(
        slide,
        9.2,
        2.0,
        3.15,
        0.4,
        [{"text": "採用理由", "size": 17, "bold": True, "font_name": FONT_HEAD, "color": WHITE}],
        margin=0.0,
    )
    add_textbox(
        slide,
        9.2,
        2.55,
        3.1,
        2.65,
        [
            {"text": "1. ローカル実行しやすい", "size": 12, "bold": True, "color": WHITE},
            {"text": "SQLite と Spring Boot を中心に構成し、デモ準備を単純化しやすい。", "size": 10.8, "color": PALE_BLUE},
            {"text": "2. 説明しやすい", "size": 12, "bold": True, "color": WHITE, "space_after": 2},
            {"text": "Web、サービス、リポジトリ、AI 連携の責務を段階的に説明できる。", "size": 10.8, "color": PALE_BLUE},
            {"text": "3. AI 連携の前提が明確", "size": 12, "bold": True, "color": WHITE, "space_after": 2},
            {"text": "プロバイダ切替、API キー、エラー応答の条件を整理しやすい。", "size": 10.8, "color": PALE_BLUE},
        ],
        margin=0.0,
    )
    add_footer_note(slide, "build.gradle と application.properties、design.md の記述をもとに整理。")


def build_environment_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, WHITE)
    add_title(slide, "6. 必要な環境と実行前提", "デモや説明に必要な前提条件を明示する", page=6)

    add_rect(slide, 0.75, 1.75, 4.25, 4.65, CARD_BG, line=LINE)
    add_badge(slide, 1.0, 2.0, 1.65, 0.34, "前提環境")
    add_textbox(
        slide,
        1.0,
        2.45,
        3.75,
        3.55,
        [
            {"text": "1. Java 21", "size": 13, "bold": True, "font_name": FONT_HEAD},
            {"text": "Gradle toolchain で Java 21 を前提にしている。", "size": 11, "color": MUTED},
            {"text": "2. ローカル実行", "size": 13, "bold": True, "font_name": FONT_HEAD, "space_after": 2},
            {"text": "`./gradlew bootRun` 実行後に `http://localhost:8080/portfolios` を開く。", "size": 11, "color": MUTED},
            {"text": "3. AI API キー", "size": 13, "bold": True, "font_name": FONT_HEAD, "space_after": 2},
            {"text": "OpenAI は `OPENAI_API_KEY`。Anthropic は `app.ai.provider=anthropic` と `ANTHROPIC_API_KEY`。", "size": 11, "color": MUTED},
            {"text": "4. ローカル保存", "size": 13, "bold": True, "font_name": FONT_HEAD, "space_after": 2},
            {"text": "SQLite ファイルは `data/portfolio-copilot.db` に作成される。", "size": 11, "color": MUTED},
        ],
        margin=0.0,
    )

    add_rect(slide, 5.3, 1.75, 3.45, 4.65, PALE_BLUE)
    add_textbox(
        slide,
        5.58,
        2.0,
        2.95,
        0.42,
        [{"text": "起動手順", "size": 16, "bold": True, "font_name": FONT_HEAD, "color": NAVY}],
        margin=0.0,
    )
    for idx, (title, body) in enumerate(
        [
            ("1. API キーを設定", "OpenAI または Anthropic の環境変数を設定する。"),
            ("2. アプリを起動", "`./gradlew bootRun` を実行する。"),
            ("3. ブラウザでアクセス", "`/portfolios` を開く。"),
            ("4. 動作確認", "相談送信とドラフト生成を試し、明示エラーも確認する。"),
        ],
        start=1,
    ):
        y = 2.55 + (idx - 1) * 0.88
        add_number_circle(slide, 5.58, y, 0.38, idx, fill=NAVY, font_size=12)
        add_textbox(
            slide,
            6.08,
            y - 0.02,
            2.2,
            0.25,
            [{"text": title, "size": 11.5, "bold": True}],
            margin=0.0,
        )
        add_textbox(
            slide,
            6.08,
            y + 0.24,
            2.25,
            0.42,
            [{"text": body, "size": 10, "color": MUTED}],
            margin=0.0,
        )

    add_rect(slide, 9.0, 1.75, 3.65, 4.65, SOFT_AMBER)
    add_textbox(
        slide,
        9.28,
        2.0,
        3.1,
        0.42,
        [{"text": "運用上の注意", "size": 16, "bold": True, "font_name": FONT_HEAD, "color": RGBColor(0x9B, 0x6B, 0x29)}],
        margin=0.0,
    )
    add_textbox(
        slide,
        9.28,
        2.55,
        3.05,
        3.3,
        [
            {"text": "設定不足のまま使わない", "size": 12, "bold": True},
            {"text": "API キー未設定時は明示的に失敗する前提。回避のための既定値投入はしない。", "size": 11, "color": MUTED},
            {"text": "出力の位置づけを保つ", "size": 12, "bold": True, "space_after": 2},
            {"text": "提案書は内部検討用ドラフトとして扱う。規制上の助言を直接確定する用途ではない。", "size": 11, "color": MUTED},
            {"text": "デモ前提を固定する", "size": 12, "bold": True, "space_after": 2},
            {"text": "API キー設定、起動手順、サンプルデータの状態を事前に確認する。", "size": 11, "color": MUTED},
        ],
        margin=0.0,
    )
    add_footer_note(slide, "必要環境は README、build.gradle、application.properties の記述に合わせている。")


def build_summary_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_background(slide, NAVY)
    add_title(slide, "7. まとめ", "説明時に押さえたい見どころ", dark=True, page=7)

    add_feature_card(
        slide,
        0.75,
        1.9,
        2.85,
        1.95,
        1,
        "機能",
        "一覧確認、保有比率、AI 相談、提案書ドラフト生成が一続きで見える。",
        RGBColor(0xD9, 0xE6, 0xF2),
    )
    add_feature_card(
        slide,
        3.92,
        1.9,
        2.85,
        1.95,
        2,
        "操作",
        "同一画面で確認から生成まで進められ、説明の流れを切りにくい。",
        RGBColor(0xD9, 0xE6, 0xF2),
    )
    add_feature_card(
        slide,
        7.09,
        1.9,
        2.85,
        1.95,
        3,
        "構成",
        "Spring Boot、SQLite、AI API の責務分離が分かりやすい。",
        RGBColor(0xD9, 0xE6, 0xF2),
    )
    add_feature_card(
        slide,
        10.26,
        1.9,
        2.35,
        1.95,
        4,
        "前提",
        "Java 21、Gradle、API キー、ブラウザが主な前提条件。",
        RGBColor(0xD9, 0xE6, 0xF2),
    )

    add_rect(slide, 0.75, 4.35, 11.86, 1.6, RGBColor(0x1F, 0x46, 0x69))
    add_textbox(
        slide,
        1.05,
        4.68,
        11.2,
        0.88,
        [
            {"text": "この資料は、実装済みのアプリ機能と仕様書をもとに構成している。説明時は、確認、相談、ドラフト生成の順で追うと全体像を共有しやすい。", "size": 16, "color": WHITE, "align": PP_ALIGN.CENTER},
        ],
        margin=0.0,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_footer_note(slide, "Portfolio Copilot の説明資料として `docs/slides/portfolio-copilot-briefing.pptx` に出力。", dark=True)


def main():
    prs = make_prs()
    build_title_slide(prs)
    build_overview_slide(prs)
    build_features_slide(prs)
    build_operation_slide(prs)
    build_architecture_slide(prs)
    build_stack_slide(prs)
    build_environment_slide(prs)
    build_summary_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"generated: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
